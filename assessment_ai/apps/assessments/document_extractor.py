import io


def extract_text_from_pdf(file_bytes: bytes) -> str:
    import fitz  # PyMuPDF
    text_parts = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def extract_text(file_bytes: bytes, filename: str) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif name.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {filename}. Use PDF, DOCX, or PPTX.")
